# -*- coding: utf-8 -*-
"""
See http://pbpython.com/creating-powerpoint.html for details on this script
Requires https://python-pptx.readthedocs.org/en/latest/index.html

Program takes a PowerPoint input file and generates a marked up version that
shows the various layouts and placeholders in the template.
"""

from __future__ import print_function
from pptx import Presentation
import argparse
import base64
import tempfile
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.table import CT_Table
from pptx.spec import GRAPHIC_DATA_URI_CHART, GRAPHIC_DATA_URI_TABLE
from pptx.util import Inches


@classmethod
def new_table_graphicFrame(cls, id_, name, rows, cols, x, y, cx, cy):
    """
    重载pptx库中graphfrm.py的该方法,以便可以指定表格的样式
    问题是这个styleId貌似对应的实际样式还是不太准确??不知道为啥了??
    """
    tableStyleId = '{5FD0F851-EC5A-4D38-B0AD-8093EC10F338}'
    graphicFrame = cls.new_graphicFrame(id_, name, x, y, cx, cy)
    graphicFrame.graphic.graphicData.uri = GRAPHIC_DATA_URI_TABLE
    graphicFrame.graphic.graphicData.append(
        CT_Table.new_tbl(rows, cols, cx, cy, tableStyleId)
    )
    return graphicFrame
CT_GraphicalObjectFrame.new_table_graphicFrame = new_table_graphicFrame


def analyze_ppt(input, output, product_dicts):
    """
    遍历输入的产品字典数据,并分析输入的方案模版,输出到ppt文件中.
    """
    prs = Presentation(input)

    for i in product_dicts:
        temp_product_obj = product_dicts[i]
        # Each powerpoint file has multiple layouts
        # 每个PowerPoint文件有多种布局(注意这里会把母版中的布局也获取出来)
        # Loop through them all and  see where the various elements are
        # 循环这些布局,并找出其中各种元素的位置
        # 这里默认母版只有一张,且为产品显示页面的布局
        for index, _ in enumerate(prs.slide_layouts):
            # 在ppt最后添加一张幻灯片,这里不知道怎么从指定位置插入,需要扩展类库原有的方法?
            slide = prs.slides.add_slide(prs.slide_layouts[index])

            # 插入产品详情的表格
            rows = 6
            cols = 2
            left = Inches(4.0)
            top = Inches(2.0)
            width = Inches(5.0)
            height = Inches(3.0)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # 设置列宽
            table.columns[0].width = Inches(1.5)
            table.columns[1].width = Inches(3.5)

            # 设置表头
            table.cell(0, 0).text = u'名称'
            table.cell(0, 1).text = temp_product_obj.name_template

            # 设置表体内容
            table.cell(1, 0).text = u'品牌'
            table.cell(1, 1).text = temp_product_obj.product_tmpl_id.prod_brand or ''
            table.cell(2, 0).text = u'型号'
            table.cell(2, 1).text = temp_product_obj.product_tmpl_id.prod_model or ''
            table.cell(3, 0).text = u'规格'
            table.cell(3, 1).text = temp_product_obj.product_tmpl_id.prod_spec or ''
            table.cell(4, 0).text = u'市场价'
            table.cell(4, 1).text = str(temp_product_obj.lst_price)
            table.cell(5, 0).text = u'产品描述'
            table.cell(5, 1).text = temp_product_obj.product_tmpl_id.description or '\n\n\n\n\n'

            # 遍历所有占位符，并通过类型识别它们
            for shape in slide.placeholders:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    try:
                        # 如果占位符为标题,设置为产品名称
                        if phf.type == PP_PLACEHOLDER.TITLE:
                            shape.text = temp_product_obj.name_template

                        # 如果占位符为图片,从数据库中读取并保存到ppt中
                        elif phf.type == PP_PLACEHOLDER.PICTURE:
                            # 如果产品有照片数据
                            if temp_product_obj.image:
                                # 先从数据库中的base64字符串中解出图片数据
                                imgdata = base64.b64decode(temp_product_obj.image)
                                # 生成一个临时文件来保存图片数据
                                temp = tempfile.TemporaryFile()
                                try:
                                    temp.write(imgdata)
                                    temp.seek(0)
                                    # 插入到图片占位符中
                                    shape.insert_picture(temp)
                                finally:
                                    # Automatically cleans up the file
                                    temp.close()
                        # 其余占位符忽略
                        else:
                            pass

                    except AttributeError:
                        print("{} has no text attribute".format(phf.type))

    prs.save(output)
